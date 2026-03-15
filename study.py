import os
import pandas as pd
from docx import Document
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, TextLoader, CSVLoader
from langchain_core.documents import Document as LC_Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_ollama import OllamaEmbeddings, ChatOllama
from langchain_community.vectorstores import Chroma
from langchain_classic.chains import RetrievalQA

def load_extra_formats(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    if ext == ".docx":
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif ext == ".xlsx":
        df = pd.read_excel(file_path)
        text = df.to_string()
    elif ext == ".pptx":
        prs = Presentation(file_path)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return [LC_Document(page_content=text, metadata={"source": file_path})] if text else []

def main():
    all_docs = []
    print("📂 Scanning folder for all study materials...")
    
    for file in os.listdir("."):
        ext = os.path.splitext(file)[1].lower()
        try:
            if ext == ".pdf":
                all_docs.extend(PyPDFLoader(file).load())
            elif ext == ".txt":
                all_docs.extend(TextLoader(file, encoding='utf-8').load())
            elif ext == ".csv":
                all_docs.extend(CSVLoader(file).load())
            elif ext in [".docx", ".xlsx", ".pptx"]:
                all_docs.extend(load_extra_formats(file))
        except Exception as e:
            print(f"⚠️ Could not read {file}: {e}")

    if not all_docs:
        print("❌ No supported files found in the folder!")
        return

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = text_splitter.split_documents(all_docs)

    print(f"🧠 Creating AI memory from {len(all_docs)} files...")
    embeddings = OllamaEmbeddings(model="nomic-embed-text")
    vector_db = Chroma.from_documents(documents=chunks, embedding=embeddings)

    llm = ChatOllama(model="tinydolphin", temperature=0.3)
    qa_chain = RetrievalQA.from_chain_type(llm=llm, retriever=vector_db.as_retriever(), chain_type="stuff")

    print("\n✅ READY! You can now ask questions about ALL your files.")
    while True:
        query = input("\nStudent Doubt: ")
        if query.lower() in ['exit', 'quit']: break
        response = qa_chain.invoke(query)
        print(f"\nAI Tutor: {response['result']}")

if __name__ == "__main__":
    main()
